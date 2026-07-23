"""
ROA OCR — Multi-Format Document Converter
===========================================
Converts DOCX, PPTX, XLSX, and HTML files to PDF for processing
through the ROA OCR pipeline. Enables multi-format document ingestion.

Supported formats:
  - DOCX (Microsoft Word) via python-docx
  - PPTX (Microsoft PowerPoint) via python-pptx
  - XLSX (Microsoft Excel) via openpyxl
  - HTML via built-in conversion
  - Plain text (.txt, .md, .csv)
"""

import io
import logging
import tempfile
import shutil
from pathlib import Path
from typing import Optional, Tuple

log = logging.getLogger("roa.converter")

# Supported input formats
SUPPORTED_FORMATS = {
    ".pdf", ".bmp", ".png", ".jpg", ".jpeg", ".tiff", ".tif",
    ".docx", ".pptx", ".xlsx", ".html", ".htm",
    ".txt", ".md", ".csv",
}

IMAGE_FORMATS = {".bmp", ".png", ".jpg", ".jpeg", ".tiff", ".tif"}
OFFICE_FORMATS = {".docx", ".pptx", ".xlsx"}
TEXT_FORMATS = {".txt", ".md", ".csv"}
WEB_FORMATS = {".html", ".htm"}


def is_supported(file_path: Path) -> bool:
    """Check if a file format is supported."""
    return file_path.suffix.lower() in SUPPORTED_FORMATS


def needs_conversion(file_path: Path) -> bool:
    """Check if a file needs conversion to PDF before OCR."""
    ext = file_path.suffix.lower()
    return ext in OFFICE_FORMATS or ext in TEXT_FORMATS or ext in WEB_FORMATS


def convert_to_pdf(src: Path, dst: Optional[Path] = None) -> Tuple[bool, Path, str]:
    """
    Convert a document to PDF format.

    Args:
        src: Source file path
        dst: Destination PDF path (auto-generated if None)

    Returns:
        (success, output_path, error_message)
    """
    ext = src.suffix.lower()

    if dst is None:
        dst = src.with_suffix(".pdf")

    if ext == ".pdf":
        if src != dst:
            shutil.copy2(src, dst)
        return True, dst, ""

    if ext in IMAGE_FORMATS:
        return _convert_image_to_pdf(src, dst)
    elif ext == ".docx":
        return _convert_docx_to_pdf(src, dst)
    elif ext == ".pptx":
        return _convert_pptx_to_pdf(src, dst)
    elif ext == ".xlsx":
        return _convert_xlsx_to_pdf(src, dst)
    elif ext in TEXT_FORMATS:
        return _convert_text_to_pdf(src, dst)
    elif ext in WEB_FORMATS:
        return _convert_html_to_pdf(src, dst)
    else:
        return False, dst, f"Unsupported format: {ext}"


def _convert_image_to_pdf(src: Path, dst: Path) -> Tuple[bool, Path, str]:
    """Convert an image to PDF."""
    try:
        from PIL import Image
        img = Image.open(src)
        if img.mode == "RGBA":
            img = img.convert("RGB")
        img.save(dst, "PDF", resolution=300)
        log.info(f"Converted image {src.name} → PDF")
        return True, dst, ""
    except Exception as e:
        return False, dst, str(e)


def _convert_docx_to_pdf(src: Path, dst: Path) -> Tuple[bool, Path, str]:
    """Convert DOCX to PDF by extracting text and rendering to PDF."""
    try:
        from docx import Document as DocxDocument
    except ImportError:
        return False, dst, "python-docx not installed. Install with: pip install python-docx"

    try:
        doc = DocxDocument(str(src))
        text_content = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                text_content.append(text)

        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                text_content.append(" | ".join(cells))
            text_content.append("")

        full_text = "\n".join(text_content)
        return _text_to_pdf(full_text, dst, title=src.stem)

    except Exception as e:
        return False, dst, f"DOCX conversion failed: {e}"


def _convert_pptx_to_pdf(src: Path, dst: Path) -> Tuple[bool, Path, str]:
    """Convert PPTX to PDF by extracting text from all slides."""
    try:
        from pptx import Presentation
    except ImportError:
        return False, dst, "python-pptx not installed. Install with: pip install python-pptx"

    try:
        prs = Presentation(str(src))
        text_content = []

        for slide_num, slide in enumerate(prs.slides, 1):
            text_content.append(f"--- Slide {slide_num} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            text_content.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        text_content.append(" | ".join(cells))
            text_content.append("")

        full_text = "\n".join(text_content)
        return _text_to_pdf(full_text, dst, title=src.stem)

    except Exception as e:
        return False, dst, f"PPTX conversion failed: {e}"


def _convert_xlsx_to_pdf(src: Path, dst: Path) -> Tuple[bool, Path, str]:
    """Convert XLSX to PDF by extracting all sheets as text."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return False, dst, "openpyxl not installed. Install with: pip install openpyxl"

    try:
        wb = load_workbook(str(src), data_only=True)
        text_content = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_content.append(f"--- Sheet: {sheet_name} ---")

            for row in ws.iter_rows(values_only=True):
                cells = [str(cell) if cell is not None else "" for cell in row]
                if any(c.strip() for c in cells):
                    text_content.append(" | ".join(cells))

            text_content.append("")

        full_text = "\n".join(text_content)
        return _text_to_pdf(full_text, dst, title=src.stem)

    except Exception as e:
        return False, dst, f"XLSX conversion failed: {e}"


def _convert_text_to_pdf(src: Path, dst: Path) -> Tuple[bool, Path, str]:
    """Convert plain text/markdown/CSV to PDF."""
    try:
        text = src.read_text(encoding="utf-8", errors="replace")
        return _text_to_pdf(text, dst, title=src.stem)
    except Exception as e:
        return False, dst, f"Text conversion failed: {e}"


def _convert_html_to_pdf(src: Path, dst: Path) -> Tuple[bool, Path, str]:
    """Convert HTML to PDF by stripping tags and rendering text."""
    try:
        import re
        html = src.read_text(encoding="utf-8", errors="replace")
        # Strip HTML tags
        text = re.sub(r'<[^>]+>', ' ', html)
        text = re.sub(r'\s+', ' ', text).strip()
        return _text_to_pdf(text, dst, title=src.stem)
    except Exception as e:
        return False, dst, f"HTML conversion failed: {e}"


def _text_to_pdf(text: str, dst: Path, title: str = "Document") -> Tuple[bool, Path, str]:
    """Render plain text content to a PDF file using reportlab or PIL fallback."""
    # Try reportlab first (best quality)
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import mm
        from reportlab.pdfgen import canvas

        dst.parent.mkdir(parents=True, exist_ok=True)
        c = canvas.Canvas(str(dst), pagesize=A4)
        width, height = A4
        margin = 25 * mm
        y = height - margin
        line_height = 14

        c.setFont("Helvetica", 10)
        for line in text.split("\n"):
            if y < margin:
                c.showPage()
                c.setFont("Helvetica", 10)
                y = height - margin

            # Truncate very long lines
            if len(line) > 100:
                words = line.split()
                current = ""
                for word in words:
                    if len(current) + len(word) + 1 > 95:
                        c.drawString(margin, y, current)
                        y -= line_height
                        current = word
                        if y < margin:
                            c.showPage()
                            c.setFont("Helvetica", 10)
                            y = height - margin
                    else:
                        current = f"{current} {word}" if current else word
                if current:
                    c.drawString(margin, y, current)
                    y -= line_height
            else:
                c.drawString(margin, y, line)
                y -= line_height

        c.save()
        log.info(f"Text rendered to PDF: {dst.name}")
        return True, dst, ""

    except ImportError:
        pass

    # Fallback: create a minimal PDF with text using PIL
    try:
        from PIL import Image, ImageDraw, ImageFont

        img_width, img_height = 2480, 3508  # A4 at 300 DPI
        margin = 150
        images = []
        current_img = Image.new("RGB", (img_width, img_height), "white")
        draw = ImageDraw.Draw(current_img)
        y = margin
        line_height = 36

        try:
            font = ImageFont.truetype("arial.ttf", 28)
        except OSError:
            font = ImageFont.load_default()

        for line in text.split("\n"):
            if y + line_height > img_height - margin:
                images.append(current_img)
                current_img = Image.new("RGB", (img_width, img_height), "white")
                draw = ImageDraw.Draw(current_img)
                y = margin

            draw.text((margin, y), line[:120], fill="black", font=font)
            y += line_height

        images.append(current_img)

        dst.parent.mkdir(parents=True, exist_ok=True)
        if len(images) == 1:
            images[0].save(dst, "PDF", resolution=300)
        else:
            images[0].save(dst, "PDF", resolution=300, save_all=True, append_images=images[1:])

        log.info(f"Text rendered to PDF via PIL: {dst.name}")
        return True, dst, ""

    except Exception as e:
        return False, dst, f"PDF rendering failed: {e}"
